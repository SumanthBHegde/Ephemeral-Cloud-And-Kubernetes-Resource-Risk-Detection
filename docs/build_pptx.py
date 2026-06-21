"""Render docs/slides.md into a Societe Generale-styled PowerPoint (docs/slides.pptx).

Usage:
    python docs/build_pptx.py            # reads docs/slides.md -> writes docs/slides.pptx

Single dependency: python-pptx. slides.md stays the single source of truth for deck content;
this script only formats it into a 16:9 widescreen deck in the Societe Generale palette
(red #E60028 / black / white) with a cover slide, per-slide kicker + headline + footer page
numbers, SG-red table headers, fitted figures, and speaker notes placed in the notes pane.

Each `---` in slides.md is a slide. A slide may contain: a `## Slide N — <kicker>` line, a
`# Big Title` (cover/close), a `### Headline` statement, bullet/numbered lists, a GitHub pipe
table, an image (![alt](figures/x.png)), `**Label:** value` meta lines (cover), and a
`*Speaker note: ...*` line. Missing image files degrade into a styled placeholder.
"""
from __future__ import annotations

import os
import re
import sys

try:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
    from pptx.util import Emu, Inches, Pt
except ImportError:  # pragma: no cover
    sys.exit("python-pptx is required: pip install python-pptx  (or: pip install -r requirements.txt)")

DOCS_DIR = os.path.dirname(os.path.abspath(__file__))
MD_PATH = os.path.join(DOCS_DIR, "slides.md")
OUT_PATH = os.path.join(DOCS_DIR, "slides.pptx")

# --- Societe Generale theme -------------------------------------------------------------------
SG_RED = RGBColor(0xE6, 0x00, 0x28)
BLACK = RGBColor(0x1D, 0x1D, 0x1B)
GRAY = RGBColor(0x60, 0x60, 0x60)
LIGHT_GRAY = RGBColor(0xB0, 0xB0, 0xB0)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BAND = RGBColor(0xF7, 0xF2, 0xF3)        # light banded-row tint
HEAD = "Calibri"
BODY = "Calibri"

EMU_W = Inches(13.333)
EMU_H = Inches(7.5)
_warnings: list[str] = []


# --- low-level shape helpers ------------------------------------------------------------------
def add_rect(slide, l, t, w, h, fill):
    from pptx.enum.shapes import MSO_SHAPE
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.fill.background()
    shp.shadow.inherit = False
    return shp


def add_text(slide, text, l, t, w, h, *, size, color=BLACK, bold=False, italic=False,
             align=PP_ALIGN.LEFT, font=BODY, anchor=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    f = run.font
    f.size = Pt(size)
    f.bold = bold
    f.italic = italic
    f.name = font
    f.color.rgb = color
    return box


def add_kicker(slide, text):
    add_text(slide, text.upper(), Inches(0.6), Inches(0.33), Inches(12.1), Inches(0.4),
             size=12, color=SG_RED, bold=True)


def add_footer(slide, page):
    add_text(slide, "Sentinel  ·  Société Générale Hackathon", Inches(0.6), Inches(7.04),
             Inches(8), Inches(0.35), size=8, color=LIGHT_GRAY)
    add_text(slide, str(page), Inches(12.4), Inches(7.04), Inches(0.5), Inches(0.35),
             size=8, color=LIGHT_GRAY, align=PP_ALIGN.RIGHT)


def fit_picture(slide, path, l, t, w, h):
    abs_path = os.path.normpath(os.path.join(DOCS_DIR, path))
    if not os.path.exists(abs_path):
        _warnings.append(f"missing image: {path}")
        add_rect(slide, l, t, w, h, BAND)
        add_text(slide, f"[ figure pending — {path} ]", l, t + h / 2 - Inches(0.2), w, Inches(0.4),
                 size=11, color=GRAY, italic=True, align=PP_ALIGN.CENTER)
        return
    pic = slide.shapes.add_picture(abs_path, l, t)
    scale = min(w / pic.width, h / pic.height)
    pic.width = int(pic.width * scale)
    pic.height = int(pic.height * scale)
    pic.left = int(l + (w - pic.width) / 2)
    pic.top = int(t + (h - pic.height) / 2)


def add_bullets(slide, items, l, t, w, h, *, size=16):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(10)
        p.line_spacing = 1.08
        dot = p.add_run()
        dot.text = "▸  "
        dot.font.size = Pt(size)
        dot.font.color.rgb = SG_RED
        dot.font.bold = True
        run = p.add_run()
        run.text = item
        run.font.size = Pt(size)
        run.font.name = BODY
        run.font.color.rgb = BLACK
    return box


def add_table(slide, rows, l, t, w, h):
    nrows, ncols = len(rows), len(rows[0])
    gf = slide.shapes.add_table(nrows, ncols, l, t, w, h)
    table = gf.table
    # column widths: first column widest
    if ncols == 3:
        widths = [0.52, 0.24, 0.24]
    elif ncols == 2:
        widths = [0.5, 0.5]
    else:
        widths = [1.0 / ncols] * ncols
    for c in range(ncols):
        table.columns[c].width = Emu(int(w * widths[c]))
    for r in range(nrows):
        for c in range(ncols):
            cell = table.cell(r, c)
            cell.margin_left = Inches(0.12)
            cell.margin_right = Inches(0.08)
            cell.margin_top = Inches(0.04)
            cell.margin_bottom = Inches(0.04)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.fill.solid()
            if r == 0:
                cell.fill.fore_color.rgb = SG_RED
            else:
                cell.fill.fore_color.rgb = WHITE if r % 2 == 1 else BAND
            tf = cell.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT if c == 0 else PP_ALIGN.CENTER
            run = p.add_run()
            run.text = rows[r][c]
            run.font.name = BODY
            run.font.size = Pt(13 if r == 0 else 12)
            run.font.bold = (r == 0)
            run.font.color.rgb = WHITE if r == 0 else BLACK
    return gf


# --- inline cleanup ---------------------------------------------------------------------------
def strip_inline(text: str) -> str:
    text = re.sub(r"\*\*(.+?)\*\*", r"\1", text)
    text = re.sub(r"`([^`]+)`", r"\1", text)
    text = re.sub(r"(?<!\w)[*_]([^*_\s][^*_]*?)[*_](?!\w)", r"\1", text)
    return text.strip()


def _split_row(line: str) -> list[str]:
    s = line.strip()
    if s.startswith("|"):
        s = s[1:]
    if s.endswith("|"):
        s = s[:-1]
    return [strip_inline(c.strip()) for c in s.split("|")]


def _is_sep(line: str) -> bool:
    return "|" in line and "-" in line and bool(re.match(r"^[\s:|-]+$", line.strip()))


# --- slide parsing ----------------------------------------------------------------------------
class Slide:
    def __init__(self):
        self.kicker = ""
        self.big_title = ""
        self.headline = ""
        self.bullets: list[str] = []
        self.table: list[list[str]] = []
        self.image: tuple[str, str] | None = None   # (alt, path)
        self.meta: list[tuple[str, str]] = []
        self.note = ""


def parse_slides(md: str) -> list[Slide]:
    chunks = re.split(r"^\s*---\s*$", md, flags=re.MULTILINE)
    slides: list[Slide] = []
    for chunk in chunks:
        if "## Slide" not in chunk:
            continue
        s = Slide()
        lines = chunk.splitlines()
        i, n = 0, len(lines)
        while i < n:
            raw = lines[i]
            line = raw.strip()
            if not line:
                i += 1
                continue
            mk = re.match(r"^##\s+Slide\s+\d+\s*[—-]\s*(.*)$", line)
            if mk:
                s.kicker = mk.group(1).strip()
                i += 1
                continue
            mb = re.match(r"^#\s+(.*)$", line)
            if mb:
                s.big_title = strip_inline(mb.group(1))
                i += 1
                continue
            mh = re.match(r"^###\s+(.*)$", line)
            if mh:
                s.headline = (s.headline + " " if s.headline else "") + strip_inline(mh.group(1))
                i += 1
                continue
            mn = re.match(r"^\*Speaker note:\s*(.*?)\*$", line)
            if mn:
                s.note = mn.group(1).strip()
                i += 1
                continue
            mi = re.match(r"^!\[(.*)\]\((.+)\)\s*$", line)
            if mi:
                s.image = (strip_inline(mi.group(1)), mi.group(2))
                i += 1
                continue
            # table
            if "|" in line and i + 1 < n and _is_sep(lines[i + 1]):
                buf = [lines[i], lines[i + 1]]
                i += 2
                while i < n and "|" in lines[i] and lines[i].strip():
                    buf.append(lines[i])
                    i += 1
                s.table = [_split_row(buf[0])] + [_split_row(r) for r in buf[2:]]
                continue
            # meta line  **Label:** value
            mm = re.match(r"^\*\*(.+?):\*\*\s*(.*)$", line)
            if mm and not re.match(r"^[-*]\s", line):
                s.meta.append((mm.group(1).strip(), strip_inline(mm.group(2).strip())))
                i += 1
                continue
            # bullets
            lb = re.match(r"^[-*]\s+(.*)$", line)
            nb = re.match(r"^\d+\.\s+(.*)$", line)
            if lb:
                s.bullets.append(strip_inline(lb.group(1)))
                i += 1
                continue
            if nb:
                s.bullets.append(strip_inline(nb.group(1)))
                i += 1
                continue
            i += 1
        slides.append(s)
    return slides


# --- slide renderers --------------------------------------------------------------------------
def _headline_size(text: str) -> int:
    n = len(text)
    if n > 95:
        return 17
    if n > 60:
        return 19
    return 22


def render_cover(slide, s: Slide):
    add_rect(slide, 0, 0, EMU_W, Inches(0.62), SG_RED)
    add_text(slide, "SOCIÉTÉ GÉNÉRALE HACKATHON  ·  CLOUD SECURITY GOVERNANCE & RISK",
             Inches(0.6), Inches(0.12), Inches(12.2), Inches(0.4),
             size=12, color=WHITE, bold=True)
    add_rect(slide, Inches(0.6), Inches(2.15), Inches(0.16), Inches(2.4), SG_RED)
    add_text(slide, s.big_title or "Sentinel", Inches(0.95), Inches(2.0), Inches(11.5), Inches(1.4),
             size=54, color=BLACK, bold=True)
    if s.headline:
        add_text(slide, s.headline, Inches(0.98), Inches(3.5), Inches(11.4), Inches(1.4),
                 size=19, color=GRAY)
    y = 5.35
    for label, val in s.meta:
        add_text(slide, f"{label}:", Inches(0.98), Inches(y), Inches(1.8), Inches(0.4),
                 size=13, color=SG_RED, bold=True)
        add_text(slide, val, Inches(2.7), Inches(y), Inches(9.5), Inches(0.4),
                 size=13, color=BLACK)
        y += 0.42


def render_close(slide, s: Slide, page):
    add_kicker(slide, s.kicker)
    add_rect(slide, Inches(0.6), Inches(2.0), Inches(12.1), Inches(0.04), SG_RED)
    add_text(slide, s.big_title, Inches(0.6), Inches(2.3), Inches(12.1), Inches(1.1),
             size=40, color=SG_RED, bold=True, align=PP_ALIGN.CENTER)
    if s.headline:
        add_text(slide, s.headline, Inches(1.2), Inches(3.6), Inches(10.9), Inches(0.9),
                 size=18, color=BLACK, align=PP_ALIGN.CENTER)
    if s.bullets:
        add_bullets(slide, s.bullets, Inches(3.2), Inches(4.7), Inches(7.0), Inches(2.0), size=14)
    add_footer(slide, page)


def render_content(slide, s: Slide, page):
    add_kicker(slide, s.kicker)
    headline = s.headline or s.kicker
    hsize = _headline_size(headline)
    add_text(slide, headline, Inches(0.6), Inches(0.72), Inches(12.1), Inches(1.05),
             size=hsize, color=BLACK, bold=True)
    add_rect(slide, Inches(0.6), Inches(1.78), Inches(12.1), Inches(0.035), SG_RED)

    body_t = Inches(2.0)
    body_h = Inches(4.85)
    has_img = s.image is not None
    has_tbl = bool(s.table)
    has_bul = bool(s.bullets)

    if has_img and (has_tbl or has_bul):
        # two columns: text/table left, figure right
        left_l, left_w = Inches(0.6), Inches(5.7)
        right_l, right_w = Inches(6.6), Inches(6.1)
        if has_tbl:
            add_table(slide, s.table, left_l, body_t, left_w, Inches(3.6))
        elif has_bul:
            add_bullets(slide, s.bullets, left_l, body_t, left_w, body_h, size=16)
        fit_picture(slide, s.image[1], right_l, body_t, right_w, body_h)
    elif has_img:
        fit_picture(slide, s.image[1], Inches(1.7), body_t, Inches(9.9), body_h)
    elif has_tbl:
        add_table(slide, s.table, Inches(1.2), body_t, Inches(10.9), Inches(3.8))
    elif has_bul:
        add_bullets(slide, s.bullets, Inches(0.9), body_t, Inches(11.5), body_h, size=18)
    add_footer(slide, page)


def main() -> int:
    if not os.path.exists(MD_PATH):
        sys.exit(f"not found: {MD_PATH}")
    with open(MD_PATH, "r", encoding="utf-8") as f:
        md = f.read()
    slides = parse_slides(md)

    prs = Presentation()
    prs.slide_width = EMU_W
    prs.slide_height = EMU_H
    blank = prs.slide_layouts[6]

    for idx, s in enumerate(slides, start=1):
        slide = prs.slides.add_slide(blank)
        kick = s.kicker.lower()
        if "title" in kick or (s.big_title and s.meta):
            render_cover(slide, s)
        elif "close" in kick or (s.big_title and not s.meta):
            render_close(slide, s, idx)
        else:
            render_content(slide, s, idx)
        if s.note:
            slide.notes_slide.notes_text_frame.text = s.note

    prs.save(OUT_PATH)
    print(f"wrote {OUT_PATH}")
    print(f"slides: {len(slides)}")
    if _warnings:
        print(f"warnings ({len(_warnings)}):", file=sys.stderr)
        for w in _warnings:
            print(f"  - {w}", file=sys.stderr)
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
